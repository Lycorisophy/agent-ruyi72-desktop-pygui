#!/usr/bin/env python3
"""
PPT Presentation Writer for OpenClaw
Creates professional .pptx presentations
"""

import argparse
import json
import sys
from pathlib import Path
from datetime import datetime


def create_ppt_document(content_dict: dict, output_path: str) -> dict:
    """
    创建 PPT 演示文稿
    
    Args:
        content_dict: 演示文稿内容配置
        output_path: 输出文件路径
    
    Returns:
        创建结果
    """
    results = {
        'success': False,
        'path': output_path,
        'message': ''
    }
    
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
    except ImportError:
        results['message'] = """
❌ 缺少依赖库 python-pptx

请安装依赖：
pip install python-pptx

或者使用 md-writer 技能创建 Markdown 演示文稿。
"""
        return results
    
    prs = Presentation()
    
    # 幻灯片布局
    layouts = {
        'title': 0,           # 标题幻灯片
        'title_content': 1,   # 标题和内容
        'section': 2,         # 分节标题
        'two_content': 3,     # 两栏内容
        'blank': 6,           # 空白
    }
    
    # 标题幻灯片
    if 'title' in content_dict:
        slide = prs.slides.add_slide(prs.slide_layouts[layouts['title']])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = content_dict['title']
        if 'subtitle' in content_dict:
            subtitle.text = content_dict['subtitle']
    
    # 幻灯片列表
    for slide_data in content_dict.get('slides', []):
        layout_name = slide_data.get('layout', 'title_content')
        layout_idx = layouts.get(layout_name, 1)
        
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        
        # 标题
        if 'title' in slide_data:
            slide.shapes.title.text = slide_data['title']
        
        # 内容
        if 'content' in slide_data:
            content = slide.placeholders[1] if slide.placeholders[1] else slide.shapes.placeholders[1]
            content.text = slide_data['content']
        
        # 列表
        if 'bullets' in slide_data:
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            for i, bullet in enumerate(slide_data['bullets']):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = bullet
                p.level = slide_data.get('level', 0)
        
        # 图片
        if 'image' in slide_data:
            try:
                slide.shapes.add_picture(
                    slide_data['image'],
                    Inches(1),
                    Inches(2),
                    width=Inches(8)
                )
            except Exception:
                pass
    
    # 保存文件
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    
    results['success'] = True
    results['message'] = f"""
✅ PPT 演示文稿创建成功！

📄 文件路径: {output_path}
📏 文件大小: {output_path.stat().st_size / 1024:.1f} KB
📊 幻灯片数量: {len(prs.slides)}

可直接在 Microsoft PowerPoint 或其他兼容软件中打开。
"""
    return results


def format_results(results: dict) -> str:
    """格式化结果"""
    return results.get('message', '❌ 创建失败')


def main():
    parser = argparse.ArgumentParser(
        description='PPT Presentation Writer for OpenClaw',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python ppt_writer.py --title "演示" --content '{"slides": [{"title": "概述", "content": "内容"}]}'
  python ppt_writer.py --output "presentation.pptx" --data '{"title": "项目汇报", "subtitle": "2026年", "slides": [{"title": "背景", "bullets": ["点1", "点2"]}]}'
        """
    )
    
    parser.add_argument('--output', '-o', default='presentation.pptx', help='Output file path')
    parser.add_argument('--title', '-t', default='新演示文稿', help='Presentation title')
    parser.add_argument('--subtitle', '-s', default='', help='Subtitle for title slide')
    parser.add_argument('--content', '-c', default='{}', help='Content as JSON')
    parser.add_argument('--data', '-d', help='Full presentation data as JSON')
    
    args = parser.parse_args()
    
    if args.data:
        content_dict = json.loads(args.data)
    else:
        content_dict = json.loads(args.content)
    
    if 'title' not in content_dict:
        content_dict['title'] = args.title
    if 'subtitle' not in content_dict and args.subtitle:
        content_dict['subtitle'] = args.subtitle
    
    results = create_ppt_document(content_dict, args.output)
    print(format_results(results))


if __name__ == '__main__':
    main()
